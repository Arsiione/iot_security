from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "IoT_Security_Scanner_prezentasiya.pptx"
SCREENSHOTS = ROOT / "screenshots"
VISUALS = ROOT / "diploma_assets" / "program_visuals_v2"

BG = RGBColor(18, 19, 31)
PANEL = RGBColor(34, 35, 53)
PANEL_2 = RGBColor(42, 43, 63)
CYAN = RGBColor(0, 194, 255)
GREEN = RGBColor(60, 235, 137)
YELLOW = RGBColor(255, 196, 77)
TEXT = RGBColor(246, 248, 255)
MUTED = RGBColor(171, 177, 196)
LINE = RGBColor(82, 86, 112)


def inch(value: float):
    return Inches(value)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int = 22,
             color: RGBColor = TEXT, bold: bool = False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Segoe UI"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, size: int = 15,
                color: RGBColor = TEXT):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def add_rect(slide, x: float, y: float, w: float, h: float, fill: RGBColor = PANEL,
             line: RGBColor = LINE, radius: bool = True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_accent(slide, x: float, y: float, w: float = 0.08, h: float = 0.54, color: RGBColor = CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_header(slide, title: str, subtitle: str | None = None, number: str | None = None):
    add_accent(slide, 0.55, 0.45, 0.08, 0.62)
    add_text(slide, title, 0.75, 0.34, 8.8, 0.45, size=25, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.76, 0.86, 9.4, 0.32, size=10, color=MUTED)
    if number:
        add_text(slide, number, 12.3, 0.35, 0.55, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def add_footer(slide):
    add_text(slide, "IoT Security Scanner  |  Lokal tor howpsuzlygy", 0.55, 7.15, 4.8, 0.22, size=8, color=MUTED)


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float, pad: float = 0.06,
                        frame: bool = True):
    if frame:
        add_rect(slide, x, y, w, h, fill=PANEL_2, line=LINE, radius=True)
    with Image.open(path) as im:
        iw, ih = im.size
    box_w = max(w - 2 * pad, 0.1)
    box_h = max(h - 2 * pad, 0.1)
    img_ratio = iw / ih
    box_ratio = box_w / box_h
    if img_ratio > box_ratio:
        pic_w = box_w
        pic_h = box_w / img_ratio
    else:
        pic_h = box_h
        pic_w = box_h * img_ratio
    px = x + (w - pic_w) / 2
    py = y + (h - pic_h) / 2
    return slide.shapes.add_picture(str(path), inch(px), inch(py), inch(pic_w), inch(pic_h))


def add_card(slide, title: str, body: str, x: float, y: float, w: float, h: float,
             accent: RGBColor = CYAN):
    add_rect(slide, x, y, w, h, fill=PANEL, line=LINE, radius=True)
    add_accent(slide, x + 0.18, y + 0.19, 0.05, 0.36, accent)
    add_text(slide, title, x + 0.34, y + 0.14, w - 0.5, 0.26, size=13, bold=True)
    add_text(slide, body, x + 0.22, y + 0.58, w - 0.42, h - 0.72, size=9, color=MUTED)


def add_slide_number(slide, idx: int):
    add_text(slide, f"{idx:02d}", 12.35, 7.07, 0.45, 0.22, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def build():
    prs = Presentation()
    prs.slide_width = inch(13.333)
    prs.slide_height = inch(7.5)
    blank = prs.slide_layouts[6]

    screenshots = {
        "home": SCREENSHOTS / "1.png",
        "scan": SCREENSHOTS / "2.png",
        "panel": SCREENSHOTS / "3.png",
        "results": SCREENSHOTS / "4.png",
        "history": SCREENSHOTS / "5.png",
    }
    visuals = {
        "arch": VISUALS / "01_programma_arhitekturasy_v2.png",
        "algo": VISUALS / "02_skanirleme_algoritmi_v2.png",
        "net": VISUALS / "03_lokal_tor_kartasy_v2.png",
        "tech": VISUALS / "04_tehnologiyalar_v2.png",
        "risk": VISUALS / "06_risk_modeli_v2.png",
    }

    # 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_text(slide, "IoT Security Scanner", 0.65, 0.65, 6.2, 0.55, size=33, bold=True)
    add_text(slide, "Lokal tordaky IoT gurluşlarynyň howpsuzlygyny awtomatlaşdyrylan barlamak", 0.68, 1.28, 5.7, 0.7, size=15, color=MUTED)
    add_bullets(slide, ["Adapter saýlamak", "Gurluşlary tapmak", "Portlary barlamak", "Risk we maslahat bermek"], 0.72, 2.28, 4.2, 1.7, size=14)
    for i, label in enumerate(["Python", "PyQt6", "Nmap", "SQLite"]):
        add_card(slide, label, "taslamanyň tehniki esasy", 0.7 + i * 1.35, 5.35, 1.18, 0.78, accent=[CYAN, GREEN, YELLOW, CYAN][i])
    add_picture_contain(slide, screenshots["scan"], 6.15, 0.72, 6.55, 5.45, pad=0.07)
    add_footer(slide); add_slide_number(slide, 1)

    # 2
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "Mesele we taslamanyň maksady", "IoT gurluşlary köpeldi, emma olaryň howpsuzlyk ýagdaýy köplenç gözegçiliksiz galýar.", "02")
    add_picture_contain(slide, visuals["net"], 7.1, 1.38, 5.55, 4.45, pad=0.05)
    add_card(slide, "Mesele", "IP kameralar, routerler, akylly öý enjamlary we telefonlar lokal torda görünýär, ýöne olaryň açyk portlary we töwekgelçiligi el bilen çalt bahalandyrylmaýar.", 0.72, 1.55, 5.55, 1.32, accent=YELLOW)
    add_card(slide, "Çözgüt", "Programma Wi‑Fi/Ethernet adapterini saýlaýar, IP aralygyny kesgitleýär, gurluşlary tapýar, portlary barlaýar we ulanyja düşnükli maslahat berýär.", 0.72, 3.05, 5.55, 1.32, accent=CYAN)
    add_card(slide, "Goragda görkezilýän netije", "Ulanyjy gara konsolsyz EXE programmasyny açýar, skanirleme geçirýär we netijeleri tablisa, taryh hem-de PDF hasabat görnüşinde görýär.", 0.72, 4.55, 5.55, 1.32, accent=GREEN)
    add_footer(slide)

    # 3
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "Ulanylan tehnologiýalar", "Taslama ýeňil, açyk çeşmeli we Windows üçin amatly tehnologiýalara esaslanýar.", "03")
    add_picture_contain(slide, visuals["tech"], 7.1, 1.25, 5.55, 4.72, pad=0.04)
    tech_cards = [
        ("Python", "programma logikasy, prosesler we faýllar"),
        ("PyQt6", "desktop interfeýs, sahypalar we tablisa"),
        ("Nmap", "host we port skanirlemesi"),
        ("SQLite", "skanirleme taryhy we maglumatlar bazasy"),
        ("ReportLab", "PDF hasabat taýýarlamak"),
        ("PyInstaller", "EXE görnüşinde ýygnamak"),
    ]
    for i, (t, b) in enumerate(tech_cards):
        x = 0.72 + (i % 2) * 3.0
        y = 1.35 + (i // 2) * 1.38
        add_card(slide, t, b, x, y, 2.72, 1.0, accent=[CYAN, GREEN, YELLOW, CYAN, GREEN, YELLOW][i])
    add_footer(slide)

    # 4
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "Programma arhitekturasy", "Her modul öz wezipesini ýerine ýetirýär we umumy skanirleme akymyna birleşýär.", "04")
    add_picture_contain(slide, visuals["arch"], 0.72, 1.25, 7.2, 4.9, pad=0.05)
    add_bullets(slide, [
        "UI moduly: sahypalar, düwmeler, tablisa we gara tema",
        "Core moduly: adapter, host tapmak, port barlamak",
        "Plugins: Telnet we IoT hyzmatlary boýunça gowşaklyk barlagy",
        "Database: netijeleri SQLite bazasynda saklamak",
        "Reports: PDF hasabat we diagrammalar",
    ], 8.28, 1.48, 4.1, 3.65, size=13)
    add_card(slide, "Üstünligi", "Modully gurluş geljekde täze plugin, täze hasabat ýa-da goşmaça skanirleme usulyny goşmagy ýeňilleşdirýär.", 8.25, 5.22, 4.15, 0.95, accent=GREEN)
    add_footer(slide)

    # 5
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "Skanirleme algoritmi", "Programma diňe port gözlemeýär, ilki gurluşyň torda bardygyny dürli usullar bilen tassyklaýar.", "05")
    add_picture_contain(slide, visuals["algo"], 6.55, 1.15, 6.05, 5.05, pad=0.04)
    steps = [
        ("1", "Adapter", "Wi‑Fi/Ethernet ileri tutulýar, VPN awtomatik saýlanmaýar."),
        ("2", "Gurluş tapmak", "Nmap, ARP cache, Ping we Gateway maglumatlary birleşdirilýär."),
        ("3", "Port barlagy", "IoT-de köp duş gelýän portlar barlanýar."),
        ("4", "Risk", "Açyk hyzmatlara görä pes, orta ýa-da ýokary töwekgelçilik berilýär."),
    ]
    for i, (n, t, b) in enumerate(steps):
        y = 1.35 + i * 1.18
        add_rect(slide, 0.75, y, 5.25, 0.92, fill=PANEL, line=LINE)
        add_text(slide, n, 0.98, y + 0.17, 0.35, 0.28, size=16, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, t, 1.45, y + 0.12, 1.75, 0.25, size=13, bold=True)
        add_text(slide, b, 1.45, y + 0.43, 4.25, 0.27, size=9, color=MUTED)
    add_footer(slide)

    # 6
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "Interfeýs: baş sahypa we skanirleme", "Ulanyjy üçin esasy ýol: programma açylýar, adapter saýlanýar we skanirleme başlanýar.", "06")
    add_picture_contain(slide, screenshots["home"], 0.72, 1.22, 5.95, 3.32, pad=0.04)
    add_picture_contain(slide, screenshots["scan"], 6.82, 1.22, 5.8, 3.32, pad=0.04)
    add_card(slide, "Baş sahypa", "Umumy ýagdaý, soňky maglumatlar we çep menýu arkaly ähli sahypalara geçiş.", 0.72, 4.85, 5.95, 0.96, accent=CYAN)
    add_card(slide, "Skanirleme sahypasy", "Adapter, IP aralygy, progress bar, netijeler tablisası we real wagt loglary.", 6.82, 4.85, 5.8, 0.96, accent=GREEN)
    add_footer(slide)

    # 7
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "Netijeleri görmek: Panel, Netijeler, Taryh", "Programma skanirleme gutarandan soň netijeleri dürli görnüşde görkezýär.", "07")
    add_picture_contain(slide, screenshots["panel"], 0.62, 1.18, 4.1, 2.6, pad=0.04)
    add_picture_contain(slide, screenshots["results"], 4.85, 1.18, 4.1, 2.6, pad=0.04)
    add_picture_contain(slide, screenshots["history"], 9.08, 1.18, 3.75, 2.6, pad=0.04)
    add_card(slide, "Panel", "Umumy statistika, soňky skanirleme we howpsuzlyk ýagdaýy.", 0.62, 4.18, 4.1, 1.05, accent=CYAN)
    add_card(slide, "Netijeler", "IP, MAC, öndüriji, görnüş, tapylan usul, portlar we töwekgelçilik.", 4.85, 4.18, 4.1, 1.05, accent=GREEN)
    add_card(slide, "Taryh", "Öňki barlaglary saklamak, deňeşdirmek we hasabat üçin ulanmak.", 9.08, 4.18, 3.75, 1.05, accent=YELLOW)
    add_footer(slide)

    # 8
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_header(slide, "Goragda görkezilýän amaly netije", "Taslama diňe nazary däl, işläp duran desktop programma görnüşinde görkezilýär.", "08")
    add_picture_contain(slide, visuals["risk"], 7.02, 1.15, 5.62, 4.78, pad=0.04)
    add_card(slide, "1. Dogry tor saýlanýar", "Programma Wi‑Fi/Ethernet adapterini tapýar we IP aralygyny awtomatik doldurýar.", 0.72, 1.32, 5.55, 0.88, accent=CYAN)
    add_card(slide, "2. Gurluşlar görkezilýär", "Router, kompýuter, telefon ýa-da IoT kandidat gurluşlar tablisa düşýär.", 0.72, 2.42, 5.55, 0.88, accent=GREEN)
    add_card(slide, "3. Töwekgelçilik düşündirilýär", "Açyk portlar we hyzmatlara görä risk derejesi hem-de maslahat berilýär.", 0.72, 3.52, 5.55, 0.88, accent=YELLOW)
    add_card(slide, "4. Hasabat taýýarlanýar", "Netijeler bazada saklanýar we diplom goragy üçin PDF hasabat eksport edilýär.", 0.72, 4.62, 5.55, 0.88, accent=CYAN)
    add_text(slide, "Netije: IoT torunyň başlangyç howpsuzlyk auditi has çalt, düşnükli we görkezmek üçin amatly bolýar.", 0.78, 6.18, 11.1, 0.38, size=15, bold=True, color=TEXT)
    add_footer(slide)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    out = build()
    print(out)
